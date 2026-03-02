"""
PPTX file parser for Microsoft PowerPoint presentations.
Uses python-pptx library for text extraction.
"""

import os
from typing import Tuple, Optional, Callable
from .base_parser import BaseParser


class PptxParser(BaseParser):
    """Parser for Microsoft PowerPoint (.pptx) files"""

    def __init__(self):
        """Initialize PPTX parser"""
        self.supported_extensions = {'.pptx'}

    def parse(self, file_path: str, cancel_check: Optional[Callable[[], bool]] = None) -> dict:
        """
        Parse PPTX file and extract text content.

        Args:
            file_path: Path to the PPTX file
            cancel_check: Optional callable that returns True if operation should be cancelled

        Returns:
            Dictionary with parsing results
        """
        # Validate file first
        is_valid, message = self.validate(file_path)
        if not is_valid:
            return {
                'success': False,
                'content': '',
                'metadata': {'file_path': file_path, 'error': message},
                'error': message
            }

        try:
            # Import pptx here to avoid dependency issues
            from pptx import Presentation

            # Load presentation
            prs = Presentation(file_path)

            # Extract text from all slides with cancellation checks
            slide_texts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                # Check cancellation before processing each slide
                if cancel_check and cancel_check():
                    return {
                        'success': True,
                        'content': '\n\n'.join(slide_texts),
                        'metadata': {'file_path': file_path, 'cancelled': True},
                        'error': None,
                        'cancelled': True
                    }

                slide_content = []

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)

                # Extract text from tables
                for shape in slide.shapes:
                    if hasattr(shape, 'has_table') and shape.has_table:
                        table_text = self._extract_table_text(shape.table)
                        if table_text:
                            slide_content.append(table_text)

                if slide_content:
                    slide_header = f"=== Slide {slide_num} ==="
                    slide_texts.append(f"{slide_header}\n" + "\n".join(slide_content))

            content = '\n\n'.join(slide_texts)

            # Get metadata
            metadata = self.get_metadata(file_path)

            return {
                'success': True,
                'content': content,
                'metadata': metadata,
                'error': None
            }

        except ImportError:
            return {
                'success': False,
                'content': '',
                'metadata': {'file_path': file_path},
                'error': "python-pptx library not installed. Install with: pip install python-pptx"
            }
        except Exception as e:
            return {
                'success': False,
                'content': '',
                'metadata': {'file_path': file_path},
                'error': f"Error parsing PPTX file: {str(e)}"
            }

    def _extract_table_text(self, table) -> str:
        """Extract text from a table shape"""
        table_text = []
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                if cell.text.strip():
                    row_text.append(cell.text.strip())
            if row_text:
                table_text.append(" | ".join(row_text))
        return "\n".join(table_text)

    def get_metadata(self, file_path: str) -> dict:
        """
        Get metadata about the PPTX file.

        Args:
            file_path: Path to the PPTX file

        Returns:
            Dictionary with file metadata
        """
        metadata = self._get_basic_metadata(file_path)

        try:
            from pptx import Presentation

            prs = Presentation(file_path)

            # Count slides and shapes
            slide_count = len(prs.slides)

            # Count shapes with text
            text_shape_count = 0
            table_count = 0
            chart_count = 0
            picture_count = 0

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_shape_count += 1
                    if hasattr(shape, 'has_table') and shape.has_table:
                        table_count += 1
                    if hasattr(shape, 'has_chart') and shape.has_chart:
                        chart_count += 1
                    if hasattr(shape, 'shape_type') and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        picture_count += 1

            # Get presentation properties
            core_properties = {}
            if hasattr(prs, 'core_properties'):
                props = prs.core_properties
                core_properties = {
                    'title': props.title,
                    'author': props.author,
                    'created': str(props.created) if props.created else None,
                    'modified': str(props.modified) if props.modified else None,
                    'last_modified_by': props.last_modified_by,
                    'revision': props.revision,
                    'category': props.category,
                    'keywords': props.keywords,
                    'comments': props.comments,
                }

            metadata.update({
                'slide_count': slide_count,
                'text_shape_count': text_shape_count,
                'table_count': table_count,
                'chart_count': chart_count,
                'picture_count': picture_count,
                'presentation_properties': core_properties,
                'slide_size': {
                    'width': prs.slide_width,
                    'height': prs.slide_height
                } if hasattr(prs, 'slide_width') and hasattr(prs, 'slide_height') else None,
            })

        except ImportError:
            metadata['error'] = "python-pptx library not installed"
        except Exception as e:
            metadata['error'] = f"Error reading PPTX metadata: {str(e)}"

        return metadata

    def validate(self, file_path: str) -> Tuple[bool, str]:
        """
        Validate if file is a valid PPTX file.

        Args:
            file_path: Path to the file

        Returns:
            Tuple of (is_valid: bool, message: str)
        """
        # Check if file exists and is readable
        file_valid, file_message = self._validate_file_exists(file_path)
        if not file_valid:
            return False, file_message

        # Check file extension
        _, ext = os.path.splitext(file_path)
        if ext.lower() not in self.supported_extensions:
            return False, f"File extension {ext} is not supported. Expected .pptx"

        # Check file signature (PPTX files are ZIP archives)
        # Only check if file is large enough
        try:
            if os.path.getsize(file_path) >= 4:
                with open(file_path, 'rb') as f:
                    header = f.read(4)
                    # PPTX files are ZIP archives starting with PK\x03\x04
                    if header != b'PK\x03\x04':
                        return False, "File does not appear to be a valid PPTX (missing ZIP header)"
        except (IOError, OSError):
            pass  # Skip signature check if we can't read file

        return True, "Valid PPTX file"